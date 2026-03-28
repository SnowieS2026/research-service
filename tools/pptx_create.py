"""
pptx_create.py -- Create PowerPoint (.pptx) presentations from structured input.
Usage:
  python pptx_create.py --output "out.pptx" --slides "Slide 1,Slide 2,Slide 3" --content "Slide 1:Title|Body text||Slide 2:Title|Bullet 1|Bullet 2"

Slide format in --content:
  Slide N:Title|text|body -> slide N with title and body
  Special prefixes: Bullet|, Number|, Image:|Link:|

Options:
  --output             Output .pptx path (required)
  --slides             Comma-separated slide titles
  --content            Pipe-delimited slide content
  --theme              Theme: default|prismal (default: prismal)
  --image              Add image: "SlideN,C:\path\to\img.png,position,width"
                       positions: top, center, bottom, bg
"""

import argparse
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)


# PRISMAL theme colours
PRISMAL_BG = RGBColor(0x0A, 0x0F, 0x1C)       # Dark navy bg
PRISMAL_ACCENT = RGBColor(0x00, 0xD4, 0xA1)   # Mint accent
PRISMAL_TEXT = RGBColor(0xFF, 0xFF, 0xFF)      # White text
PRISMAL_SECONDARY = RGBColor(0x9B, 0x9B, 0x9B) # Gray text


def set_slide_background(slide, color):
    """Set solid background colour for a slide."""
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs, title, subtitle=None):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, PRISMAL_BG)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = PRISMAL_ACCENT
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = PRISMAL_TEXT
        p2.alignment = PP_ALIGN.CENTER


def add_content_slide(prs, title, bullets, theme='prismal'):
    """Add a content slide with title and bullet points."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    if theme == 'prismal':
        set_slide_background(slide, PRISMAL_BG)
        title_color = PRISMAL_ACCENT
        text_color = PRISMAL_TEXT
    else:
        title_color = RGBColor(0, 0, 0)
        text_color = RGBColor(0, 0, 0)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = title_color

    # Underline accent
    line_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(2), Inches(0.05))
    line_tf = line_box.text_frame
    line_p = line_tf.paragraphs[0]
    line_p.text = ''
    from pptx.util import Pt as P2
    from pptx.enum.shapes import MSO_SHAPE
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(2), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = PRISMAL_ACCENT
    line.line.fill.background()

    # Content
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = text_color
        p.level = 0
        # Bullet marker
        pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        from lxml import etree
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '•')
        buFont = etree.SubElement(pPr, qn('a:buFont'))
        buFont.set('typeface', 'Arial')

    return slide


def add_image_to_slide(slide, image_path, position='center', width=5):
    """Add an image to a slide at a given position."""
    path = Path(image_path)
    if not path.exists():
        print(f"  WARNING: Image not found: {image_path}")
        return

    if position == 'center':
        left = (prs.slide_width - Inches(width)) / 2
        top = (prs.slide_height - Inches(width * 0.6)) / 2
    elif position == 'top':
        left = Inches(1)
        top = Inches(1.5)
    elif position == 'bottom':
        left = Inches(1)
        top = Inches(5)
    else:
        left = Inches(1)
        top = Inches(1.5)

    try:
        slide.shapes.add_picture(str(path), left, top, width=Inches(width))
    except Exception as e:
        print(f"  WARNING: Could not add image: {e}")


def process_content(content_str):
    """
    Parse content string into slides.
    Slides separated by '||'
    Each slide: "SlideTitle|body1|body2|..."
    Special: Bullet:|Number: prefix for list items
    """
    slides = []
    raw_slides = content_str.split('||')
    for raw in raw_slides:
        raw = raw.strip()
        if not raw:
            continue
        # Check for "Slide N:" or "SlideTitle:" prefix
        parts = raw.split('|')
        title = parts[0]
        bullets = parts[1:] if len(parts) > 1 else []
        slides.append((title, bullets))
    return slides


def main():
    parser = argparse.ArgumentParser(description='Create PowerPoint presentations')
    parser.add_argument('--output', '-o', required=True, help='Output .pptx path')
    parser.add_argument('--slides', '-s', help='Comma-separated slide titles (alternative to --content)')
    parser.add_argument('--content', '-c', help='Slide content with || separators')
    parser.add_argument('--theme', default='prismal', choices=['default', 'prismal'], help='Visual theme')
    parser.add_argument('--image', action='append', help='Add image: "SlideIndex,path,position,width"')
    parser.add_argument('--author', help='Author name (document properties)')
    args = parser.parse_args()

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    if args.content:
        slides = process_content(args.content)
        for title, bullets in slides:
            if not bullets and title:
                # Title-only slide
                add_title_slide(prs, title)
            else:
                add_content_slide(prs, title, bullets, theme=args.theme)

    elif args.slides:
        slide_titles = [t.strip() for t in args.slides.split(',')]
        for i, title in enumerate(slide_titles):
            if i == 0 and title.lower() in ('title', 'cover'):
                add_title_slide(prs, title)
            else:
                add_content_slide(prs, title, [], theme=args.theme)

    # Add images
    if args.image:
        for img_spec in args.image:
            parts = [p.strip() for p in img_spec.split(',')]
            if len(parts) < 2:
                continue
            slide_idx = int(parts[0])
            img_path = parts[1]
            position = parts[2] if len(parts) > 2 else 'center'
            width = float(parts[3]) if len(parts) > 3 else 5.0

            if slide_idx < len(prs.slides):
                add_image_to_slide(prs.slides[slide_idx], img_path, position, width)

    output = Path(args.output)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    print(f"Created: {output} ({output.stat().st_size:,} bytes)")


if __name__ == '__main__':
    main()
